"""
PowerPoint Generator Service
Adapted from final_year_project-main/ppt_generator.py
Creates beautiful presentations using python-pptx with multiple themes
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from typing import Dict, List, Optional
import os
import io
import logging
import textwrap

logger = logging.getLogger(__name__)


class PPTService:
    """Service for generating PowerPoint presentations"""
    
    # Color schemes for different themes
    THEMES = {
        'modern': {
            'bg_color': RGBColor(255, 255, 255),
            'title_color': RGBColor(31, 78, 121),
            'text_color': RGBColor(64, 64, 64),
            'accent_color': RGBColor(0, 120, 215)
        },
        'dark': {
            'bg_color': RGBColor(30, 30, 30),
            'title_color': RGBColor(255, 255, 255),
            'text_color': RGBColor(220, 220, 220),
            'accent_color': RGBColor(0, 120, 215)
        },
        'professional': {
            'bg_color': RGBColor(255, 255, 255),
            'title_color': RGBColor(68, 84, 106),
            'text_color': RGBColor(89, 89, 89),
            'accent_color': RGBColor(192, 0, 0)
        },
        'business': {
            'bg_color': RGBColor(248, 249, 250),
            'title_color': RGBColor(33, 37, 41),
            'text_color': RGBColor(73, 80, 87),
            'accent_color': RGBColor(0, 123, 255)
        },
        'academic': {
            'bg_color': RGBColor(255, 255, 255),
            'title_color': RGBColor(52, 58, 64),
            'text_color': RGBColor(73, 80, 87),
            'accent_color': RGBColor(111, 66, 193)
        },
        'minimal': {
            'bg_color': RGBColor(255, 255, 255),
            'title_color': RGBColor(0, 0, 0),
            'text_color': RGBColor(100, 100, 100),
            'accent_color': RGBColor(128, 128, 128)
        },
        'creative': {
            'bg_color': RGBColor(255, 250, 240),
            'title_color': RGBColor(220, 53, 69),
            'text_color': RGBColor(102, 102, 102),
            'accent_color': RGBColor(255, 193, 7)
        }
    }
    
    def __init__(self, theme: str = 'modern', brand_colors: Optional[List[str]] = None):
        """
        Initialize PPT Service
        
        Args:
            theme: Theme name (modern, dark, professional, etc.)
            brand_colors: Optional list of hex colors to override theme
        """
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        # Get theme or default to modern
        self.theme = self.THEMES.get(theme, self.THEMES['modern'])
        
        # Apply brand colors if provided
        if brand_colors:
            self._apply_brand_colors(brand_colors)
        
        logger.info(f"Initialized PPT service with theme: {theme}")
    
    def _apply_brand_colors(self, brand_colors: List[str]):
        """Apply custom brand colors to theme"""
        if len(brand_colors) >= 1:
            self.theme['accent_color'] = self._hex_to_rgb(brand_colors[0])
        if len(brand_colors) >= 2:
            self.theme['title_color'] = self._hex_to_rgb(brand_colors[1])
    
    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color to RGBColor"""
        hex_color = hex_color.lstrip('#')
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )
    
    def _wrap_text(self, text: str, width: int = 80) -> str:
        """Wrap text to prevent overflow"""
        # Handle empty or None text
        if not text:
            return text
        
        # Split by newlines first to preserve intentional line breaks
        lines = text.split('\n')
        wrapped_lines = []
        
        for line in lines:
            if len(line) <= width:
                wrapped_lines.append(line)
            else:
                # Wrap long lines
                wrapped = textwrap.fill(line, width=width, break_long_words=False, break_on_hyphens=True)
                wrapped_lines.append(wrapped)
        
        return '\n'.join(wrapped_lines)
    
    def _calculate_optimal_font_size(self, text_frame, initial_size: int = 18, min_size: int = 10) -> int:
        """Calculate optimal font size to fit text in frame"""
        current_size = initial_size
        
        # Try reducing font size if text overflows
        while current_size >= min_size:
            # Set all paragraphs to current size
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(current_size)
            
            # Check if text fits (approximation based on character count and dimensions)
            # This is a heuristic since python-pptx doesn't provide direct overflow detection
            total_chars = sum(len(p.text) for p in text_frame.paragraphs)
            frame_area = text_frame._element.getparent().getparent()
            
            # Rough estimation: if we have too many characters, reduce font size
            # Approximate: 80 chars per line at 18pt, scale accordingly
            estimated_lines = total_chars / (80 * (current_size / 18))
            
            # If estimated lines fit in available height, we're good
            if estimated_lines <= 15:  # Reasonable number of lines for a slide
                break
                
            current_size -= 1
        
        return current_size
    
    def create_title_slide(self, title: str, subtitle: str = ""):
        """
        Create a title slide
        
        Args:
            title: Main title
            subtitle: Subtitle or author info
        """
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme['bg_color']
        
        # Add title
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(1.5)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        
        # Wrap title text if too long
        wrapped_title = self._wrap_text(title, width=50)
        title_frame.text = wrapped_title
        
        # Format title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = self.theme['title_color']
        
        # Add subtitle if provided
        if subtitle:
            left = Inches(1)
            top = Inches(4.5)
            width = Inches(8)
            height = Inches(0.75)
            
            subtitle_box = slide.shapes.add_textbox(left, top, width, height)
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.word_wrap = True
            
            # Wrap subtitle text if too long
            wrapped_subtitle = self._wrap_text(subtitle, width=70)
            subtitle_frame.text = wrapped_subtitle
            
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.alignment = PP_ALIGN.CENTER
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = self.theme['text_color']
        
        logger.debug("Created title slide")
        return slide
    
    def create_content_slide(
        self, 
        title: str, 
        content: List[str], 
        image_path: Optional[str] = None,
        speaker_notes: str = ""
    ):
        """
        Create a content slide with title, bullet points, and optional image
        
        Args:
            title: Slide title
            content: List of bullet points
            image_path: Optional path to image file
            speaker_notes: Optional speaker notes
        """
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme['bg_color']
        
        # Add title
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(0.75)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        
        # Wrap title text if too long
        wrapped_title = self._wrap_text(title, width=60)
        title_frame.text = wrapped_title
        
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.theme['title_color']
        
        # Add accent line under title
        line = slide.shapes.add_shape(
            1,  # Line shape
            left=Inches(0.5),
            top=Inches(1.4),
            width=Inches(9),
            height=Inches(0)
        )
        line.line.color.rgb = self.theme['accent_color']
        line.line.width = Pt(3)
        
        # Determine layout based on whether there's an image
        if image_path and os.path.exists(image_path):
            # Content on left, image on right
            content_left = Inches(0.5)
            content_top = Inches(1.7)
            content_width = Inches(5)
            content_height = Inches(5)
            
            # Add image
            img_left = Inches(6)
            img_top = Inches(2)
            img_width = Inches(3.5)
            
            try:
                slide.shapes.add_picture(
                    image_path,
                    img_left,
                    img_top,
                    width=img_width
                )
                logger.debug(f"Added image to slide: {image_path}")
            except Exception as e:
                logger.error(f"Error adding image: {e}")
        else:
            # Full width for content
            content_left = Inches(0.5)
            content_top = Inches(1.7)
            content_width = Inches(9)
            content_height = Inches(5)
        
        # Add content
        content_box = slide.shapes.add_textbox(
            content_left, content_top, content_width, content_height
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.vertical_anchor = MSO_ANCHOR.TOP
        
        # Determine appropriate wrap width based on content width
        wrap_width = int(content_width.inches * 12)  # Approximate chars per inch
        
        for i, bullet in enumerate(content):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            # Wrap bullet text to prevent overflow
            wrapped_bullet = self._wrap_text(bullet, width=wrap_width)
            p.text = wrapped_bullet
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = self.theme['text_color']
            p.space_before = Pt(12)
        
        # Calculate and apply optimal font size if needed
        total_chars = sum(len(bullet) for bullet in content)
        num_bullets = len(content)
        
        # If content is too long, adjust font size
        if num_bullets > 8 or total_chars > 800:
            optimal_size = self._calculate_optimal_font_size(content_frame, initial_size=18, min_size=12)
            logger.debug(f"Adjusted content font size to {optimal_size}pt to prevent overflow")
        
        # Add speaker notes if provided
        if speaker_notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = speaker_notes
            logger.debug("Added speaker notes to slide")
        
        logger.debug(f"Created content slide: {title}")
        return slide
    
    def create_chart_slide(
        self,
        title: str,
        chart_data: Dict,
        chart_type: str = "bar"
    ):
        """
        Create a slide with a chart
        
        Args:
            title: Slide title
            chart_data: Dictionary with 'categories' and 'values' lists
            chart_type: Type of chart (bar, line, pie)
        """
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme['bg_color']
        
        # Add title
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(0.75)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        
        # Wrap title text if too long
        wrapped_title = self._wrap_text(title, width=60)
        title_frame.text = wrapped_title
        
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.theme['title_color']
        
        # Prepare chart data
        chart_data_obj = CategoryChartData()
        chart_data_obj.categories = chart_data.get('categories', [])
        chart_data_obj.add_series('Series 1', chart_data.get('values', []))
        
        # Determine chart type
        chart_type_map = {
            'bar': XL_CHART_TYPE.BAR_CLUSTERED,
            'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'line': XL_CHART_TYPE.LINE,
            'pie': XL_CHART_TYPE.PIE
        }
        
        xl_chart_type = chart_type_map.get(chart_type, XL_CHART_TYPE.BAR_CLUSTERED)
        
        # Add chart
        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(5)
        chart = slide.shapes.add_chart(
            xl_chart_type, x, y, cx, cy, chart_data_obj
        ).chart
        
        logger.debug(f"Created chart slide: {title}")
        return slide
    
    def create_thank_you_slide(self):
        """Create a thank you / closing slide"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme['bg_color']
        
        # Add "Thank You" text
        left = Inches(1)
        top = Inches(3)
        width = Inches(8)
        height = Inches(1.5)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.text = "Thank You!"
        
        para = text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(60)
        para.font.bold = True
        para.font.color.rgb = self.theme['title_color']
        
        logger.debug("Created thank you slide")
        return slide
    
    def generate_from_data(
        self, 
        presentation_data: Dict, 
        image_paths: Optional[Dict[int, str]] = None
    ) -> bytes:
        """
        Generate complete presentation from data and return as bytes
        
        Args:
            presentation_data: Dictionary with presentation structure
            image_paths: Dictionary mapping slide numbers to image paths
            
        Returns:
            PPTX file as bytes
        """
        image_paths = image_paths or {}
        
        # Create title slide
        title = presentation_data.get('title', 'Presentation')
        subtitle = presentation_data.get('subtitle', '')
        self.create_title_slide(title, subtitle)
        
        # Create content slides
        slides = presentation_data.get('slides', [])
        
        for slide_data in slides:
            slide_num = slide_data.get('slide_number') or slide_data.get('index', 0) + 1
            slide_title = slide_data.get('title', f'Slide {slide_num}')
            slide_content = slide_data.get('content') or slide_data.get('bullets', [])
            speaker_notes = slide_data.get('speaker_notes', '')
            image_path = image_paths.get(slide_num)
            
            # Check if this should be a chart slide
            if slide_data.get('chart_data'):
                self.create_chart_slide(
                    slide_title,
                    slide_data['chart_data'],
                    slide_data.get('chart_type', 'bar')
                )
            else:
                self.create_content_slide(
                    slide_title,
                    slide_content,
                    image_path,
                    speaker_notes
                )
        
        # Create thank you slide
        self.create_thank_you_slide()
        
        # Save to bytes
        ppt_bytes = io.BytesIO()
        self.prs.save(ppt_bytes)
        ppt_bytes.seek(0)
        
        logger.info(f"Generated presentation with {len(self.prs.slides)} slides")
        
        return ppt_bytes.getvalue()
